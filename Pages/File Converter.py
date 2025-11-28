#!/usr/bin/env python

"""
An advanced, interactive file converter for Termux (Text UI Version).
Version 3.2: Organizes all folders into a master "File Converter" directory.

This script will:
1. Check and install required Python libraries.
2. Check for external binaries (ffmpeg, unrar, cairo).
3. Create a "File Converter" folder in Downloads, containing 40 organizer sub-folders.
4. Provide a numbered text-based UI for navigation.
5. Support many document, image, A/V, and data conversions.

CRITICAL SETUP (in Termux):
pkg install ffmpeg unrar libcairo libgirepository libjpeg-turbo libpng
"""

import sys
import os
import subprocess
import importlib.util
import time
# import curses # Removed for Text UI
import traceback
import zipfile
import tarfile
import csv
import json
import gzip
import shutil
from contextlib import redirect_stderr, redirect_stdout

# --- 1. SETUP & CONFIGURATION ---

# (14) Python libraries to auto-install
REQUIRED_MODULES = {
    "Pillow": "Pillow",         # Images
    "reportlab": "reportlab",   # PDF creation
    "docx": "python-docx",    # Word docs
    "odf": "odfpy",           # OpenOffice docs
    "bs4": "beautifulsoup4",  # HTML/XML parsing
    "markdown": "Markdown",     # Markdown parsing
    "lxml": "lxml",           # XML/HTML parsing
    "cairosvg": "cairosvg",     # SVG conversion
    "psd_tools": "psd-tools",   # PSD reading
    "striprtf": "striprtf",     # RTF reading
    "EbookLib": "EbookLib",     # EPUB reading
    "pptx": "python-pptx",    # PowerPoint reading
    "rarfile": "rarfile",       # RAR extraction
    "py7zr": "py7zr"          # 7-Zip extraction
}

# (40) Folders to create
FOLDER_NAMES = [
    # Images (10)
    "JPG", "PNG", "WEBP", "BMP", "TIFF", "GIF", "ICO", "TGA", "SVG", "PSD",
    # Documents (12)
    "PDF", "TXT", "DOCX", "ODT", "HTML", "MD", "CSV", "RTF", "EPUB", "JSON", "XML", "PPTX",
    # Archives (5)
    "ZIP", "TAR", "RAR", "7Z", "GZ",
    # Audio (7)
    "MP3", "WAV", "OGG", "FLAC", "M4A", "AAC", "WMA",
    # Video (6)
    "MP4", "MKV", "AVI", "MOV", "WMV", "FLV"
]

# Helper lists for logic
IMAGE_FOLDERS = ["JPG", "PNG", "WEBP", "BMP", "TIFF", "GIF", "ICO", "TGA", "SVG", "PSD"]
IMAGE_EXTS = ['.jpg', '.jpeg', '.png', '.webp', '.bmp', '.tiff', '.gif', '.ico', '.tga']
VECTOR_IMAGE_EXTS = ['.svg']
LAYERED_IMAGE_EXTS = ['.psd']
AV_EXTS = ['.mp3', '.wav', '.ogg', '.flac', '.m4a', '.aac', '.wma', '.mp4', '.mkv', '.avi', '.mov', '.wmv', '.flv']
ARCHIVE_EXTS = ['.zip', '.tar', '.gz', '.bz2', '.rar', '.7z']
TEXT_DOC_EXTS = ['.txt', '.docx', '.odt', '.html', '.md', '.csv', '.rtf', '.epub', '.json', '.xml', '.pptx', '.svg']
DATA_EXTS = ['.csv', '.json', '.xml']

# Paths
STORAGE_PATH = "/storage/emulated/0"
DOWNLOAD_PATH = os.path.join(STORAGE_PATH, "Download")
BASE_CONVERTER_PATH = os.path.join(DOWNLOAD_PATH, "File Converter")

# Global flags for external binaries
HAS_FFMPEG = False
HAS_UNRAR = False
HAS_CAIRO = False

# --- 2. PRE-UI SETUP FUNCTIONS ---

def clear_screen():
    os.system('cls' if os.name == 'nt' else 'clear')

def check_and_install_dependencies():
    """Checks and installs required Python modules."""
    print("--- Checking Required Python Libraries (14) ---")
    all_installed = True
    for module_name, package_name in REQUIRED_MODULES.items():
        spec = importlib.util.find_spec(module_name)
        if spec is None:
            all_installed = False
            print(f"Installing '{package_name}'...")
            try:
                with open(os.devnull, 'w') as devnull:
                    with redirect_stdout(devnull), redirect_stderr(devnull):
                        subprocess.run([sys.executable, "-m", "pip", "install", package_name], check=True)
                print(f"Successfully installed '{package_name}'.")
            except Exception:
                print(f"ERROR: Failed to install '{package_name}'.")
                print(f"Please install it manually: pip install {package_name}")
                sys.exit(1)
        else:
            pass
    
    if all_installed:
        print("All Python libraries are present.\n")
    else:
        print("All required libraries are now installed.\n")
    time.sleep(0.5)

def check_external_bins():
    """Checks for 'ffmpeg', 'unrar', and 'cairo'."""
    global HAS_FFMPEG, HAS_UNRAR, HAS_CAIRO
    print("--- Checking External Binaries ---")
    
    # Check ffmpeg
    try:
        with open(os.devnull, 'w') as devnull:
            subprocess.run(["ffmpeg", "-version"], check=True, stdout=devnull, stderr=devnull)
        print("Found 'ffmpeg'. Audio/Video conversions are ENABLED.")
        HAS_FFMPEG = True
    except (subprocess.CalledProcessError, FileNotFoundError):
        print("WARNING: 'ffmpeg' not found. A/V conversions DISABLED.")
        print("  To enable, run: pkg install ffmpeg\n")
        HAS_FFMPEG = False

    # Check unrar
    try:
        with open(os.devnull, 'w') as devnull:
            subprocess.run(["unrar"], check=True, stdout=devnull, stderr=devnull)
        print("Found 'unrar'. RAR extraction is ENABLED.")
        HAS_UNRAR = True
    except (subprocess.CalledProcessError, FileNotFoundError):
        print("WARNING: 'unrar' not found. RAR extraction DISABLED.")
        print("  To enable, run: pkg install unrar\n")
        HAS_UNRAR = False
        
    # Check cairo (for SVG)
    if importlib.util.find_spec("cairosvg") is not None:
        print("Found 'cairosvg'. SVG conversions are ENABLED.")
        HAS_CAIRO = True
    else:
        print("WARNING: 'cairosvg' Python lib not found. SVG conversion DISABLED.")
        print("  The script tried to install it, but it may have failed.")
        print("  You may also need: pkg install libcairo libgirepository\n")
        HAS_CAIRO = False
        
    print("")
    time.sleep(0.5)


def check_storage_access():
    print("--- Checking Storage Access ---")
    if not os.path.exists(DOWNLOAD_PATH):
        print(f"ERROR: Cannot access internal storage at '{DOWNLOAD_PATH}'.")
        print("Please run 'termux-setup-storage' in your Termux terminal,")
        print("grant the permission, and then run this script again.")
        sys.exit(1)
    print("Storage access confirmed.\n")
    time.sleep(0.5)

def setup_folders():
    print(f"--- Setting up Organizer Folders ---")
    print(f"Location: {BASE_CONVERTER_PATH}")
    try:
        os.makedirs(BASE_CONVERTER_PATH, exist_ok=True)
        for folder in FOLDER_NAMES:
            os.makedirs(os.path.join(BASE_CONVERTER_PATH, folder), exist_ok=True)
        print(f"Successfully created/verified {len(FOLDER_NAMES)} sub-folders.\n")
    except Exception as e:
        print(f"ERROR: Could not create folders: {e}")
        sys.exit(1)
    time.sleep(0.5)

# --- 3. IMPORTS (Post-Installation) ---
try:
    from PIL import Image
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from docx import Document
    from odf.opendocument import load as odf_load
    from odf.text import P as odf_P
    from bs4 import BeautifulSoup
    import markdown
    import lxml
    import cairosvg
    from psd_tools import PSDImage
    from striprtf.striprtf import rtf_to_text
    from ebooklib import epub, ITEM_DOCUMENT
    import pptx
    import rarfile
    import py7zr
except ImportError as e:
    print(f"CRITICAL ERROR: Failed to import a library: {e}")
    print("Please ensure all dependencies are installed (see startup logs).")
    sys.exit(1)

# --- 4. CORE CONVERSION LOGIC ---
# (Updated to remove stdscr arguments)

def get_text_from_file(input_path, in_ext):
    """Helper to extract raw text from various document types."""
    text_lines = []
    try:
        if in_ext == '.txt':
            with open(input_path, 'r', encoding='utf-8') as f:
                text_lines = f.readlines()
        elif in_ext == '.docx':
            doc = Document(input_path)
            text_lines = [para.text + '\n' for para in doc.paragraphs]
        elif in_ext == '.odt':
            doc = odf_load(input_path)
            for para in doc.getElementsByType(odf_P):
                text_lines.append(str(para) + '\n')
        elif in_ext in ['.html', '.xml', '.svg']:
            with open(input_path, 'r', encoding='utf-8') as f:
                parser = 'lxml' if in_ext != '.html' else 'html.parser'
                soup = BeautifulSoup(f, parser)
                text_lines = [line + '\n' for line in soup.get_text().splitlines()]
        elif in_ext == '.md':
            with open(input_path, 'r', encoding='utf-8') as f:
                html = markdown.markdown(f.read())
                soup = BeautifulSoup(html, 'html.parser')
                text_lines = [line + '\n' for line in soup.get_text().splitlines()]
        elif in_ext == '.csv':
            with open(input_path, 'r', encoding='utf-8', newline='') as f:
                reader = csv.reader(f)
                text_lines = [','.join(row) + '\n' for row in reader]
        elif in_ext == '.json':
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                text_lines = [json.dumps(data, indent=2)]
        elif in_ext == '.rtf':
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            text_lines = [rtf_to_text(content)]
        elif in_ext == '.epub':
            book = epub.read_epub(input_path)
            for item in book.get_items():
                if item.get_type() == ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text_lines.append(soup.get_text() + '\n\n') # Add space between chapters
        elif in_ext == '.pptx':
            prs = pptx.Presentation(input_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_lines.append(shape.text + '\n')
    except Exception as e:
        raise Exception(f"Text extraction failed: {e}")
    return text_lines

def write_text_to_pdf(text_lines, output_path):
    c = canvas.Canvas(output_path, pagesize=A4)
    width, height = A4
    margin_x, margin_y = 0.75 * inch, 1 * inch
    text_object = c.beginText(margin_x, height - margin_y)
    text_object.setFont("Helvetica", 10)
    line_height, y = 12, height - margin_y
    for line in text_lines:
        for sub_line in line.split('\n'):
            if y < margin_y:
                c.drawText(text_object)
                c.showPage()
                text_object = c.beginText(margin_x, height - margin_y)
                text_object.setFont("Helvetica", 10)
                y = height - margin_y
            text_object.textLine(sub_line.strip('\r'))
            y -= line_height
    c.drawText(text_object)
    c.save()

def handle_image_conversion(in_path, out_path):
    with Image.open(in_path) as img:
        if out_path.lower().endswith(('.jpg', '.jpeg')):
            if img.mode == 'RGBA':
                img = img.convert('RGB')
        img.save(out_path)

def handle_svg_conversion(in_path, out_path):
    if not HAS_CAIRO:
        raise Exception("Cairo/SVG libraries not installed.")
    out_ext = os.path.splitext(out_path)[1].lower()
    if out_ext == '.png':
        cairosvg.svg2png(url=in_path, write_to=out_path)
    elif out_ext == '.pdf':
        cairosvg.svg2pdf(url=in_path, write_to=out_path)
    else:
        raise Exception(f"SVG conversion to {out_ext} not supported.")

def handle_psd_conversion(in_path, out_path):
    psd = PSDImage.open(in_path)
    composite_image = psd.composite()
    composite_image.save(out_path)

def handle_av_conversion(in_path, out_path):
    if not HAS_FFMPEG:
        raise Exception("'ffmpeg' not found. A/V conversion is disabled.")
    command = ['ffmpeg', '-i', in_path, '-y', out_path]
    
    print("--- Running ffmpeg ---")
    print(f"Command: {' '.join(command)}")
    print("This may take some time...")
    try:
        with open(os.devnull, 'w') as devnull:
            subprocess.run(command, check=True, stdout=devnull, stderr=subprocess.STDOUT)
        print("ffmpeg conversion finished successfully.")
    except Exception as e:
        print(f"ffmpeg ERROR: {e}")
        raise Exception(f"ffmpeg conversion failed. {e}")

def handle_extraction(in_path, out_folder_path, in_ext):
    base_name = os.path.splitext(os.path.basename(in_path))[0]
    extract_path = os.path.join(out_folder_path, base_name)
    os.makedirs(extract_path, exist_ok=True)
    
    if in_ext == '.zip':
        with zipfile.ZipFile(in_path, 'r') as zf:
            zf.extractall(extract_path)
    elif in_ext in ['.tar', '.gz', '.bz2']:
        if in_ext == '.gz' and not in_path.endswith('.tar.gz'):
             out_filename = os.path.splitext(os.path.basename(in_path))[0]
             out_path = os.path.join(out_folder_path, out_filename)
             with gzip.open(in_path, 'rb') as f_in:
                 with open(out_path, 'wb') as f_out:
                     shutil.copyfileobj(f_in, f_out)
             return f"Decompressed to: {out_path}"
        else:
            with tarfile.open(in_path, 'r:*') as tf:
                tf.extractall(extract_path)
    elif in_ext == '.rar':
        if not HAS_UNRAR:
            raise Exception("'unrar' binary not found.")
        with rarfile.RarFile(in_path) as rf:
            rf.extractall(extract_path)
    elif in_ext == '.7z':
        with py7zr.SevenZipFile(in_path, 'r') as zf:
            zf.extractall(extract_path)
            
    return f"Extracted to: {extract_path}"

def handle_data_conversion(in_path, out_path, in_ext, out_ext):
    if in_ext == '.csv' and out_ext == '.json':
        data = []
        with open(in_path, 'r', encoding='utf-8', newline='') as f:
            reader = csv.DictReader(f)
            for row in reader:
                data.append(row)
        with open(out_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=4)
    elif in_ext == '.json' and out_ext == '.csv':
        with open(in_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        if not isinstance(data, list) or not data:
            raise Exception("JSON must be a non-empty list of objects.")
        if not all(isinstance(x, dict) for x in data):
            raise Exception("JSON must be a list of objects (dicts).")
            
        headers = data[0].keys()
        with open(out_path, 'w', encoding='utf-8', newline='') as f:
            writer = csv.DictWriter(f, fieldnames=headers)
            writer.writeheader()
            writer.writerows(data)
    else:
        raise Exception(f"Data conversion {in_ext} to {out_ext} not supported.")

def handle_md_to_html(in_path, out_path):
    with open(in_path, 'r', encoding='utf-8') as f:
        html = markdown.markdown(f.read())
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(html)

def handle_single_image_to_pdf(in_path, out_path):
    try:
        with Image.open(in_path) as img:
            img_rgb = img.convert('RGB')
            img_rgb.save(out_path, "PDF", resolution=100.0)
    except Exception as e:
        raise Exception(f"Pillow (Image->PDF) Error: {e}")

def handle_multi_image_to_pdf(image_paths, out_path):
    try:
        images_rgb = []
        for path in image_paths:
            img = Image.open(path)
            images_rgb.append(img.convert('RGB'))
        if not images_rgb:
            raise Exception("No images found to convert.")
        images_rgb[0].save(
            out_path, "PDF", resolution=100.0,
            save_all=True, append_images=images_rgb[1:]
        )
    except Exception as e:
        raise Exception(f"Pillow (Multi-Image->PDF) Error: {e}")

# --- 5. MAIN CONVERSION ROUTER ---

def convert_file(in_path, out_folder_name):
    """
    Main router function for dispatching conversion tasks.
    Returns (success_bool, message_string)
    """
    in_ext = os.path.splitext(in_path)[1].lower()
    out_ext = f".{out_folder_name.lower()}"
    
    base_name = os.path.splitext(os.path.basename(in_path))[0]
    out_folder_path = os.path.join(BASE_CONVERTER_PATH, out_folder_name)
    out_path = os.path.join(out_folder_path, f"{base_name}{out_ext}")

    try:
        # --- Route 1: Extraction ---
        if in_ext in ARCHIVE_EXTS:
            out_folder = out_folder_path if in_ext == '.gz' else os.path.join(BASE_CONVERTER_PATH, out_folder_name)
            message = handle_extraction(in_path, out_folder, in_ext)
            return (True, message)

        # --- Route 2: SVG Conversion (to PNG, PDF) ---
        if in_ext in VECTOR_IMAGE_EXTS and out_ext in ['.png', '.pdf']:
            handle_svg_conversion(in_path, out_path)
            return (True, f"Saved to: {out_path}")

        # --- Route 3: PSD Conversion (to flat image) ---
        if in_ext in LAYERED_IMAGE_EXTS and out_ext in IMAGE_EXTS:
            handle_psd_conversion(in_path, out_path)
            return (True, f"Saved to: {out_path}")

        # --- Route 4: Image-to-Image (Pillow) ---
        if in_ext in IMAGE_EXTS and out_ext in IMAGE_EXTS:
            handle_image_conversion(in_path, out_path)
            return (True, f"Saved to: {out_path}")
            
        # --- Route 5: Single Image-to-PDF ---
        if in_ext in IMAGE_EXTS and out_ext == '.pdf':
            handle_single_image_to_pdf(in_path, out_path)
            return (True, f"Saved to: {out_path}")

        # --- Route 6: A/V-to-A/V (ffmpeg) ---
        if in_ext in AV_EXTS and out_ext in AV_EXTS:
            handle_av_conversion(in_path, out_path)
            return (True, f"Saved to: {out_path}")
            
        # --- Route 7: Data Conversion (CSV <-> JSON) ---
        if in_ext in ['.csv', '.json'] and out_ext in ['.csv', '.json']:
            handle_data_conversion(in_path, out_path, in_ext, out_ext)
            return (True, f"Saved to: {out_path}")

        # --- Route 8: MD-to-HTML ---
        if in_ext == '.md' and out_ext == '.html':
            handle_md_to_html(in_path, out_path)
            return (True, f"Saved to: {out_path}")

        # --- Route 9: Anything-to-TXT ---
        if out_ext == '.txt' and in_ext in TEXT_DOC_EXTS:
            text_lines = get_text_from_file(in_path, in_ext)
            with open(out_path, 'w', encoding='utf-8') as f:
                f.writelines(text_lines)
            return (True, f"Saved to: {out_path}")
            
        # --- Route 10: Anything-to-PDF ---
        if out_ext == '.pdf' and in_ext in TEXT_DOC_EXTS:
            text_lines = get_text_from_file(in_path, in_ext)
            write_text_to_pdf(text_lines, out_path)
            return (True, f"Saved to: {out_path}")

        # --- No Route Found ---
        return (False, f"Unsupported conversion: {in_ext} to {out_ext}")

    except Exception as e:
        return (False, f"ERROR: {str(e)}")


# --- 6. TEXT UI HELPER FUNCTIONS ---

def print_header(title):
    print("=" * 40)
    print(f"      Termux File Converter")
    print("=" * 40)
    print(f" {title}")
    print("-" * 40)

def press_enter():
    input("\nPress Enter to continue...")

def get_input(prompt):
    return input(f"{prompt}: ").strip()

def show_menu(title, options, back_text="Go Back"):
    """
    Displays a numbered menu and returns the selected item (str) or None if back.
    """
    while True:
        clear_screen()
        print_header(title)
        
        # Display options
        for i, option in enumerate(options):
            print(f" [{i+1}] {option}")
        
        print("-" * 40)
        print(f" [0] {back_text}")
        print("-" * 40)
        
        choice = input("Select an option: ").strip()
        
        if choice == '0':
            return None
        
        if choice.isdigit():
            idx = int(choice) - 1
            if 0 <= idx < len(options):
                return options[idx]
        
        print("Invalid selection. Try again.")
        time.sleep(1)

def run_multi_image_to_pdf_wizard(input_folder_path, input_folder_name):
    try:
        image_paths = [
            os.path.join(input_folder_path, f) 
            for f in os.listdir(input_folder_path) 
            if os.path.splitext(f)[1].lower() in IMAGE_EXTS
        ]
        image_paths.sort()
    except Exception as e:
        print(f"Error reading images: {e}")
        press_enter()
        return

    if not image_paths:
        print("No images found in that folder.")
        press_enter()
        return

    print(f"\nFound {len(image_paths)} images in '{input_folder_name}'.")
    confirm = input("Combine ALL these images into one PDF? (y/n): ").lower()
    if confirm != 'y':
        return

    default_name = f"{input_folder_name}_Album"
    filename = input(f"Enter PDF filename (default: {default_name}): ").strip()
    if not filename:
        filename = default_name

    out_folder_path = os.path.join(BASE_CONVERTER_PATH, "PDF")
    out_path = os.path.join(out_folder_path, f"{filename}.pdf")
    
    print("\nWorking... combining images into PDF...")
    try:
        handle_multi_image_to_pdf(image_paths, out_path)
        print(f"Success! Saved to: {out_path}")
    except Exception as e:
        print(f"ERROR: {e}")
    press_enter()

def run_help():
    clear_screen()
    print_header("Help / How to Use")
    print("This converter uses a simple 3-step process:")
    print("")
    print("1. MOVE YOUR FILES:")
    print("   Use your phone's File Manager.")
    print("   Go to: /Download/File Converter/")
    print("   Move files into the correct folder (e.g., move")
    print("   'report.docx' into the 'DOCX' folder).")
    print("")
    print("2. RUN THIS CONVERTER:")
    print("   Select 'Convert a File' from the main menu.")
    print("")
    print("3. FOLLOW THE STEPS:")
    print("   Step 1: Choose the INPUT folder (e.g., 'DOCX').")
    print("   Step 2: Choose the file you want to convert.")
    print("   Step 3: Choose the OUTPUT format (e.g., 'PDF').")
    print("")
    print("** SPECIAL CONVERSIONS **")
    print(" - Archives (ZIP, RAR, etc): Extract automatically.")
    print(" - Multi-Image PDF: Choose 'JPG' -> '[ Convert ALL to PDF ]'")
    print(" - Data: Convert CSV <-> JSON.")
    press_enter()

# --- 7. MAIN TEXT APPLICATION ---

def main():
    while True:
        main_choice = show_menu("Main Menu", ["Convert a File", "Help / How to Use"], back_text="Exit")
        
        if main_choice is None: # Exit
            print("Exiting...")
            break
        
        if main_choice == "Help / How to Use":
            run_help()
            continue
            
        # Step 1: Input Folder
        input_folder = show_menu("Step 1: Choose INPUT Folder", FOLDER_NAMES)
        if input_folder is None: continue
        
        input_folder_path = os.path.join(BASE_CONVERTER_PATH, input_folder)
        
        # Check files
        try:
            files = [f for f in os.listdir(input_folder_path) if os.path.isfile(os.path.join(input_folder_path, f))]
            files.sort()
        except Exception as e:
            print(f"Error accessing folder: {e}")
            press_enter()
            continue
            
        file_options = []
        if input_folder in IMAGE_FOLDERS and len(files) > 0:
            file_options.append("[ ** Convert ALL Images to one PDF ** ]")
        file_options.extend(files)
        
        if not file_options:
            print(f"\nNo files found in {input_folder} folder.")
            print("Please move files there using your file manager.")
            press_enter()
            continue
            
        # Step 2: Choose File
        input_file = show_menu(f"Step 2: Choose File from '{input_folder}'", file_options)
        if input_file is None: continue
        
        # Special: Multi-Image
        if input_file == "[ ** Convert ALL Images to one PDF ** ]":
            run_multi_image_to_pdf_wizard(input_folder_path, input_folder)
            continue
            
        full_input_path = os.path.join(input_folder_path, input_file)
        in_ext = os.path.splitext(input_file)[1].lower()
        
        # Step 3: Output Format / Logic
        output_folder = None
        
        if in_ext in ARCHIVE_EXTS:
            # Archives extract to self
            output_folder = input_folder
            print(f"\nSelected Archive: {input_file}")
            confirm = input(f"Extract contents? (y/n): ").lower()
        else:
            output_folder = show_menu("Step 3: Choose OUTPUT Format", FOLDER_NAMES)
            if output_folder is None: continue
            
            if output_folder == input_folder:
                print("\nError: Input and Output folder cannot be the same.")
                press_enter()
                continue
                
            print(f"\nConversion: {input_file} -> {output_folder}")
            confirm = input("Start conversion? (y/n): ").lower()
            
        if confirm == 'y':
            print("\nWorking... please wait...")
            success, message = convert_file(full_input_path, output_folder)
            if success:
                print(f"\n[SUCCESS] {message}")
            else:
                print(f"\n[FAILED] {message}")
            press_enter()

# --- 8. SCRIPT ENTRYPOINT ---

if __name__ == "__main__":
    try:
        clear_screen()
        print("--- Initializing Termux Converter v3.2 (Text UI) ---")
        check_and_install_dependencies()
        check_external_bins()
        check_storage_access()
        setup_folders()
        
        print("--- Setup Complete ---")
        print(f"Organizer folders: /storage/emulated/0/Download/File Converter/")
        print("\nStarting application...")
        time.sleep(1)
        
        main()
        
    except KeyboardInterrupt:
        print("\nExiting...")
    except Exception as e:
        print("\nA critical error occurred:")
        traceback.print_exc()